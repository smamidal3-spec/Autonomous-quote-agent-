from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BK  = RGBColor(0x00,0x00,0x00)
DK  = RGBColor(0x08,0x08,0x08)
D2  = RGBColor(0x0F,0x0F,0x0F)
D3  = RGBColor(0x18,0x18,0x18)
MID = RGBColor(0x33,0x33,0x33)
DIM = RGBColor(0x66,0x66,0x66)
LG  = RGBColor(0x99,0x99,0x99)
WH  = RGBColor(0xFF,0xFF,0xFF)
W2  = RGBColor(0xCC,0xCC,0xCC)
W3  = RGBColor(0xAA,0xAA,0xAA)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLK = prs.slide_layouts[6]

def sl():
    s = prs.slides.add_slide(BLK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BK
    return s

def box(s, l,t,w,h, fc=None, lc=None, lw=Pt(1)):
    sh = s.shapes.add_shape(1, Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid() if fc else sh.fill.background()
    if fc: sh.fill.fore_color.rgb = fc
    sh.line.color.rgb = lc if lc else (fc or BK)
    if not lc: sh.line.fill.background()
    if lc: sh.line.width = lw
    return sh

def tx(s, text, l,t,w,h, sz=Pt(14), bold=False, color=WH,
       align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.bold = bold; r.font.italic = italic
    r.font.size = sz;   r.font.color.rgb = color
    return tb

def rule(s, t, col=MID, thick=0.02):
    box(s, 0.5, t, 12.33, thick, fc=col)

def tag(s, text, t=0.38):
    tx(s, text.upper(), 0.6, t, 12.0, 0.35,
       sz=Pt(8), color=DIM, align=PP_ALIGN.CENTER)

def hdr(s, num, title):
    box(s, 0,0,13.33,0.85, fc=D2)
    tx(s, title, 0.6,0.1,11.5,0.65, sz=Pt(22), bold=True, color=WH)
    tx(s, str(num), 12.6,0.15,0.55,0.55, sz=Pt(11), color=DIM, align=PP_ALIGN.RIGHT)
    rule(s, 0.85, col=MID)

def bullets(s, pts, l,t,w,step=0.62, sz=Pt(13), col=LG):
    for i,p in enumerate(pts):
        tx(s, f"—  {p}", l, t+i*step, w, step+0.05, sz=sz, color=col)

def card(s, l,t,w,h, title, val, desc, border=MID):
    box(s, l,t,w,h, fc=D3, lc=border, lw=Pt(0.75))
    tx(s, title.upper(), l+0.15,t+0.12,w-0.3,0.3, sz=Pt(7), color=DIM)
    tx(s, val, l+0.12,t+0.42,w-0.25,0.6, sz=Pt(26), bold=True, color=WH)
    tx(s, desc, l+0.12,t+0.95,w-0.25,h-1.1, sz=Pt(11), color=LG)

def two_col(s, left_title, left_pts, right_title, right_pts):
    box(s, 0.4,1.05,6.2,5.5, fc=D2, lc=MID, lw=Pt(0.5))
    tx(s, left_title, 0.6,1.15,5.8,0.5, sz=Pt(16), bold=True, color=WH)
    rule(s, 1.65, col=D3)
    for i,p in enumerate(left_pts):
        tx(s, f"  {p}", 0.6,1.75+i*0.65,5.8,0.6, sz=Pt(12), color=LG)
    box(s, 6.72,1.05,6.2,5.5, fc=D3, lc=MID, lw=Pt(0.5))
    tx(s, right_title, 6.9,1.15,5.8,0.5, sz=Pt(16), bold=True, color=WH)
    rule(s, 1.65, col=D2)
    for i,p in enumerate(right_pts):
        tx(s, f"  {p}", 6.9,1.75+i*0.65,5.8,0.6, sz=Pt(12), color=LG)

def flow_boxes(s, labels, y=2.0, h=2.8):
    n = len(labels); w = 12.33/n - 0.02
    for i,(lbl,sub) in enumerate(labels):
        x = 0.5 + i*(w+0.02)
        box(s, x,y,w,h, fc=D3, lc=MID, lw=Pt(0.5))
        tx(s, lbl, x,y+0.25,w,0.6, sz=Pt(18), bold=True, color=WH, align=PP_ALIGN.CENTER)
        tx(s, sub, x,y+0.95,w,h-1.0, sz=Pt(10), color=LG, align=PP_ALIGN.CENTER)

def agent_triple(s, inp, logic, out):
    for i,(title,pts,col) in enumerate([(("INPUTS",inp,D2),("LOGIC",logic,D3),("OUTPUTS",out,D2))[i] for i in range(3)]):
        title2,pts2,fc = title,pts,col
        x = 0.4+i*4.32
        box(s,x,1.05,4.2,5.6, fc=fc, lc=MID, lw=Pt(0.75))
        tx(s,title2,x+0.15,1.12,4.0,0.42, sz=Pt(9), color=DIM, bold=True)
        for j,p in enumerate(pts2):
            tx(s,f"  {p}",x+0.15,1.65+j*0.68,4.0,0.62, sz=Pt(12), color=W2)

def agent_triple_fix(s, inp_pts, logic_pts, out_pts):
    panels = [("INPUTS", inp_pts, D2), ("LOGIC",  logic_pts, D3), ("OUTPUTS", out_pts, D2)]
    for i, (pname, pts, fc) in enumerate(panels):
        x = 0.4 + i*4.32
        box(s, x, 1.05, 4.2, 5.6, fc=fc, lc=MID, lw=Pt(0.75))
        tx(s, pname, x+0.15, 1.12, 4.0, 0.42, sz=Pt(9), color=DIM, bold=True)
        for j, p in enumerate(pts):
            tx(s, f"  {p}", x+0.15, 1.65+j*0.68, 4.0, 0.62, sz=Pt(12), color=W2)

S1 = sl()
box(S1,0,0,13.33,0.04, fc=WH)
box(S1,0,7.46,13.33,0.04, fc=WH)
box(S1,0,3.2,13.33,0.02, fc=D3)
tx(S1,"AUTONOMOUS INSURANCE QUOTE INTELLIGENCE",
   0.6,0.8,12.0,1.8, sz=Pt(38), bold=True, color=WH, align=PP_ALIGN.CENTER)
tx(S1,"From AI Foundations to a 6-Agent Decision System",
   0.6,2.6,12.0,0.7, sz=Pt(18), color=LG, align=PP_ALIGN.CENTER)
tx(S1,"A Cinematic Briefing for Executive Leadership  ·  March 2026",
   0.6,3.5,12.0,0.5, sz=Pt(11), color=DIM, align=PP_ALIGN.CENTER)
tx(S1,"[Company Logo]", 0.5,6.8,3.0,0.5, sz=Pt(11), color=DIM)
tx(S1,"Strictly Confidential",10.0,6.8,3.0,0.5, sz=Pt(11), color=DIM, align=PP_ALIGN.RIGHT)

S2 = sl(); hdr(S2,2,"Why This Matters")
for i,(v,lbl,d) in enumerate([
    ("68%","Loss Ratio Pressure","Industry loss ratio threatening profitability across motor lines"),
    ("₹12Cr","Annual Fraud Exposure","Undetected fraud in quote-to-bind workflows each year"),
    ("−41%","Conversion Gap","Quotes lost to misaligned pricing and absence of personalisation"),
    ("5.2d","Underwriting Cycle","Average manual turnaround — digital customers will not wait"),
]):
    card(S2, 0.4+i*3.25,1.0,3.0,2.2, lbl, v, d)
bullets(S2,[
    "Traditional rule engines cannot scale with the volume or complexity of modern data",
    "AI agents replace sequential manual steps with parallel, real-time, explainable decisions",
    "Every second saved in decisioning directly improves customer experience and conversion rate",
    "Fraud caught at quote time prevents claim losses — not after the money is gone",
],0.5,3.55,12.0)

S3 = sl(); hdr(S3,3,"The Quote Journey — Then vs Now")
two_col(S3,"Traditional",
    ["Manual form intake by the agent","Static rule lookups — 3 to 5 day cycle","Reactive fraud review after claim",
     "One-size premium — no personalisation","Decision opaque to the customer","Follow-up by phone — high drop-off"],
    "AI-Assisted",
    ["Digital API intake — validated instantly","Six agents run in under two seconds","Fraud scored at quote time — not claim time",
     "Segment-matched personalised plan offered","Every decision includes a written reason","Auto-approve, follow-up, or escalate clearly"])

S4 = sl(); hdr(S4,4,"How AI Works — Plain English")
flow_boxes(S4,[
    ("1\nData","Historical quotes, policies, claims"),
    ("2\nPatterns","Model learns what drives risk and conversion"),
    ("3\nPrediction","Each new quote scored in real time"),
    ("4\nDecision","Route to approve, follow-up, or escalate"),
    ("5\nFeedback","Outcomes retrain the model monthly"),
])
bullets(S4,[
    "The model never guesses — it learns from thousands of past cases",
    "SHAP values explain every prediction in plain English to auditors and regulators",
    "Feedback loop: underwriter decisions flow back to improve the model next cycle",
],0.5,5.2,12.0)

S5 = sl(); hdr(S5,5,"Supervised vs Unsupervised Learning")
two_col(S5,"Supervised Learning",
    ["Trained on labelled outcomes — past quotes tagged HIGH / MED / LOW",
     "Model learns: age + accidents + vehicle cost → risk tier",
     "Used for Risk Profiler, Conversion Predictor, Premium Advisor",
     "Human-verified outcomes improve accuracy with every cycle",
     "Accuracy validated on 20% held-out test set before deployment"],
    "Unsupervised Learning",
    ["No labels needed — discovers hidden patterns automatically",
     "Isolation Forest identifies anomalies that deviate from the norm",
     "KMeans clusters customers into five behavioural segments",
     "Used for Fraud Detection and the Personalisation Agent",
     "Adapts continuously as customer behaviour evolves"])

S6 = sl(); hdr(S6,6,"Our Data Ecosystem")
for i,(src,desc) in enumerate([
    ("CRM System","Customer & historical policy records"),
    ("Quote Portal","Real-time intake via API"),
    ("Claims Database","Historical loss and fraud outcomes"),
    ("Market Data","Industry benchmarks and pricing indices"),
]):
    x = 0.4+i*3.25
    box(S6,x,1.0,3.0,1.1, fc=D3, lc=MID, lw=Pt(0.5))
    tx(S6,src,x+0.12,1.08,2.8,0.42, sz=Pt(13), bold=True, color=WH)
    tx(S6,desc,x+0.12,1.5,2.8,0.45, sz=Pt(11), color=LG)
for i,(nm,d) in enumerate([
    ("RAW ZONE","Ingested CSV and API dumps — untouched, append-only, time-stamped, full audit history"),
    ("CLEANED ZONE","Null removal · Type coercion · Outlier handling · Enum enforcement · Deduplication"),
    ("FEATURE LAYER","Encoded categoricals · Engineered ratios · Normalised numerics · Risk composites"),
    ("MODEL ARTIFACTS","Serialised .pkl models + label encoders + scalers — versioned, zero-downtime swap"),
    ("API OUTPUT LAYER","JSON predictions served via FastAPI — consumed by downstream agents in real time"),
]):
    y = 2.3+i*0.92
    box(S6,0.4,y,12.5,0.82, fc=D3 if i%2==0 else D2, lc=MID, lw=Pt(0.5))
    tx(S6,nm,0.6,y+0.18,3.0,0.5, sz=Pt(9), color=DIM, bold=True)
    tx(S6,d,3.8,y+0.2,9.0,0.5, sz=Pt(12), color=LG)

S7 = sl(); hdr(S7,7,"Data Storage Architecture")
layers = [
    ("RAW ZONE","Ingested CSV + API payloads — untouched, append-only. Full audit history preserved."),
    ("CLEANED ZONE","Null removal · Type coercion · Outlier handling · Canonical Enum enforcement"),
    ("FEATURE LAYER","Encoded categoricals · Engineered risk ratios · Normalised numeric features"),
    ("MODEL ARTIFACTS","Serialised .pkl models + LabelEncoders + scalers — versioned, hot-swappable"),
    ("API OUTPUT LAYER","JSON predictions served via FastAPI with sub-200ms response time"),
]
for i,(nm,d) in enumerate(layers):
    y=1.0+i*1.12
    fc = D2 if i%2==0 else D3
    box(S7,0.4,y,12.5,1.0, fc=fc, lc=MID, lw=Pt(0.5))
    tx(S7,nm,0.6,y+0.22,3.5,0.6, sz=Pt(10), bold=True, color=W3)
    tx(S7,d,4.3,y+0.25,8.5,0.55, sz=Pt(13), color=LG)
    if i<4:
        tx(S7,"↓",6.5,y+0.98,0.5,0.2, sz=Pt(10), color=DIM, align=PP_ALIGN.CENTER)

S8 = sl(); hdr(S8,8,"Training Pipeline — End to End")
flow_boxes(S8,[
    ("Ingest","Raw CSV"),("Clean","Validate"),("Encode","Enums"),
    ("Split","80 / 20"),("Train","Models"),("Validate","Accuracy"),("Deploy",".pkl"),
], y=1.1, h=2.0)
bullets(S8,[
    "Date columns removed — Policy_Bind_DT excluded to eliminate 100% conversion model leakage",
    "LabelEncoder saves exact string maps — canonical Enum mismatches raise hard errors at inference",
    "Risk model: Random Forest — 90%+ accuracy on held-out 20% test set",
    "Conversion model: XGBoost — AUC 0.82 after complete leakage elimination",
    "Models serialised with joblib — zero-downtime hot-swap on FastAPI reload",
],0.5,3.35,12.0)

S9 = sl(); hdr(S9,9,"Model Stack Overview")
models = [
    ("Risk Profiler","Random Forest","Ensemble of trees. Robust to outliers. SHAP explainability."),
    ("Conversion Predictor","XGBoost","Gradient boosted trees. Class imbalance handled. AUC 0.82."),
    ("Fraud Detection","Isolation Forest","Unsupervised anomaly scoring on 12 engineered features."),
    ("Premium Advisor","Rule Engine","Deterministic salary × risk × premium alignment. Auditable."),
    ("Personalisation","KMeans Clustering","Five behavioural segments mapped to plan tiers."),
    ("Decision Router","Rule Engine","Synthesises all agent scores into one final action."),
]
for i,(agent,model,desc) in enumerate(models):
    r,c = divmod(i,3)
    x=0.4+c*4.32; y=1.05+r*2.7
    box(S9,x,y,4.1,2.5, fc=D3 if c%2==0 else D2, lc=MID, lw=Pt(0.75))
    tx(S9,agent,x+0.15,y+0.12,3.8,0.4, sz=Pt(10), color=DIM)
    tx(S9,model,x+0.15,y+0.52,3.8,0.55, sz=Pt(15), bold=True, color=WH)
    tx(S9,desc,x+0.15,y+1.1,3.8,1.2, sz=Pt(11), color=LG)

S10 = sl(); hdr(S10,10,"6-Agent Orchestration Architecture")
tx(S10,"Quote request enters → agents execute in sequence → router outputs decision in < 2 seconds",
   0.6,0.95,12.0,0.45, sz=Pt(12), color=DIM, align=PP_ALIGN.CENTER)
agents=[("Input\nRequest","23 fields\nvia API"),("A1\nRisk\nProfiler","risk_score\nrisk_tier"),
        ("A2\nFraud\nDetection","fraud_score\nrule_flags"),("A3\nConversion\nPredictor","probability\nband"),
        ("A4\nPerson-\nalisation","segment\nplan"),("A5\nPremium\nAdvisor","premium\nrange"),
        ("A6\nDecision\nRouter","AUTO\nFOLLOW\nESCALATE")]
for i,(nm,out) in enumerate(agents):
    x=0.35+i*1.83
    box(S10,x,1.5,1.7,2.8, fc=D3 if i%2==0 else D2, lc=MID, lw=Pt(0.75))
    tx(S10,nm,x,1.62,1.7,1.2, sz=Pt(11), bold=True, color=WH, align=PP_ALIGN.CENTER)
    tx(S10,out,x,2.85,1.7,1.3, sz=Pt(9), color=LG, align=PP_ALIGN.CENTER)
for i,(lbl,d) in enumerate([("AUTO-APPROVE","LOW risk + HIGH conversion + no flags"),
                              ("FOLLOW-UP","MEDIUM risk OR LOW conversion OR premium issue"),
                              ("ESCALATE","HIGH risk OR fraud flag OR any rule flag present")]):
    y=4.55+i*0.88
    box(S10,0.4,y,12.5,0.78, fc=D3 if i%2==0 else D2, lc=MID, lw=Pt(0.5))
    tx(S10,lbl,0.6,y+0.2,3.5,0.45, sz=Pt(12), bold=True, color=WH)
    tx(S10,d,4.3,y+0.22,8.5,0.42, sz=Pt(12), color=LG)

def make_agent(num, name, inp, logic, out, note):
    s = sl(); hdr(s, num, f"Agent {num-10}: {name}")
    agent_triple_fix(s, inp, logic, out)
    tx(s, note, 0.5, 6.82, 12.0, 0.5, sz=Pt(9), color=DIM, italic=True)
    return s

make_agent(11,"Risk Profiler",
    ["Driver Age","Driving Experience","Prior Accidents","Citation Count","Vehicle Cost Range","Coverage Type"],
    ["Random Forest — 146K clean policy records","Heuristic overrides for extreme age/experience","SHAP top-3 risk driver extraction per quote","Score range 0–100 mapped to tier","MEDIUM+ auto-triggers review queue"],
    ["risk_score (0 – 100)","risk_tier: LOW / MED / HIGH","top_3_SHAP_drivers","heuristic_override: True/False"],
    "Inputs → inputs panel left  |  Logic → centre  |  risk_score and tier → right output")

make_agent(12,"Fraud Detection",
    ["Driver Age vs Experience delta","Vehicle cost vs Salary ratio","Accident velocity","Premium vs Risk gap","Re-quote flag"],
    ["Validation rules fire first: age ≥ 18, exp ≤ age−16","Isolation Forest anomaly score on 12 features","Rule engine flags underwriting anomalies separately","fraud_flag ONLY if anomaly score ≥ 0.5","rule_flags trigger escalation — NOT fraud label"],
    ["fraud_score (0.0 – 1.0)","fraud_flag: True / False","rule_flags: [list]","fraud_reason_codes: [list]"],
    "Fraud and rule flags are separated by design — different impacts on routing and fairness")

make_agent(13,"Conversion Predictor",
    ["All 18 quote features","Risk tier from Agent 1","Coverage choice","Annual Miles Range","Salary Range"],
    ["XGBoost classifier — AUC 0.82","Policy_Bind_DT excluded — leakage eliminated","Output: bind probability 0.0 – 1.0","Banded to HIGH / MEDIUM / LOW","SHAP top conversion drivers extracted"],
    ["conversion_probability","conversion_band: H / M / L","top_conversion_drivers (SHAP)","discount_recommended flag"],
    "LOW conversion triggers premium discount recommendation — directly feeds Agent 5 logic")

make_agent(14,"Personalisation Agent",
    ["Salary Range","Vehicle Cost Range","Annual Miles Range","Coverage preference","Risk tier","Driver Age"],
    ["KMeans clustering on 5 behavioural dimensions","Customer mapped to nearest segment centroid","Segments: Budget / Standard / Premium / Fleet / Senior","Add-on rules applied per segment","Young drivers auto-receive personal accident cover"],
    ["customer_segment","recommended_plan","add_ons: [list]","segment_confidence (0.0 – 1.0)"],
    "Expensive vehicles auto-triggered for zero-depreciation and engine protection add-ons")

make_agent(15,"Premium Advisor",
    ["Quoted Premium (from form)","Risk score from Agent 1","Conversion probability","Salary Range"],
    ["Calculates expected premium from risk × salary matrix","Compares quoted vs expected range","Salary ≤ ₹25L AND quoted > 2× expected → flag","Discount coupon logic for borderline cases","Generates recommended premium min and max band"],
    ["premium_issue: True / False","recommended_range [min, max]","recommendation_reason","discount_applicable flag"],
    "Prevents under-pricing of high-risk profiles and over-pricing of low-risk customers")

make_agent(16,"Decision Router",
    ["risk_tier","fraud_flag","rule_flags","conversion_band","premium_issue flag"],
    ["AUTO-APPROVE: LOW risk + HIGH conv + no flags","ESCALATE: HIGH risk OR fraud_flag OR rule_flag","FOLLOW-UP: MED risk OR LOW conv OR premium issue","Confidence score 0.0 – 1.0 attached","Written explanation for every decision path"],
    ["decision: AUTO / FOLLOW / ESCALATE","confidence_score (0.0 – 1.0)","decision_explanation","escalation_required: True/False"],
    "Fully deterministic — every routing path is auditable, no black box — IRDAI compliant")

S17 = sl(); hdr(S17,17,"Full Quote Walkthrough — Sample Profile")
box(S17,0.4,0.95,12.5,0.8, fc=D3, lc=MID, lw=Pt(0.5))
tx(S17,"Profile: Age 24 · Male · 2 yrs experience · 2 accidents · XUV 500 ₹32 Lakh · Enhanced cover · Salary ₹25 Lakh · Quoted Premium ₹1,200",
   0.6,1.05,12.2,0.6, sz=Pt(12), color=LG)
agent_outs=[("A1 RISK","Score: 78\nTier: HIGH\nAge · Accidents · Cost"),
            ("A2 FRAUD","Score: 0.31\nFraud: False\nRule: age_vehicle"),
            ("A3 CONV.","Prob: 0.61\nBand: MED\nDriver: Coverage"),
            ("A4 PERSON.","Seg: Premium\nPlan: Enhanced+\nAdd: ZD · PA"),
            ("A5 PREMIUM","Issue: True\n₹1,400 – 1,800\nToo Low"),
            ("A6 ROUTER","ESCALATE\nConf: 0.89\nFull reason")]
for i,(nm,body) in enumerate(agent_outs):
    x=0.4+i*2.12
    box(S17,x,2.0,2.0,2.5, fc=D3 if i%2==0 else D2, lc=MID, lw=Pt(0.5))
    tx(S17,nm,x+0.1,2.08,1.8,0.4, sz=Pt(9), bold=True, color=WH)
    tx(S17,body,x+0.1,2.55,1.8,1.8, sz=Pt(11), color=LG)
box(S17,0.4,4.7,12.5,1.0, fc=D2, lc=MID, lw=Pt(0.5))
tx(S17,"FINAL DECISION: ESCALATE TO UNDERWRITER",3.0,4.82,8.0,0.4, sz=Pt(14), bold=True, color=WH, align=PP_ALIGN.CENTER)
tx(S17,"HIGH risk tier + underwriting rule flag (age/vehicle mismatch) + quoted premium below recommended range",
   0.6,5.2,12.0,0.38, sz=Pt(11), color=DIM, align=PP_ALIGN.CENTER)

S18 = sl(); hdr(S18,18,"Monitoring, Governance & Retraining")
pillars=[("DRIFT DETECTION","PSI tests run weekly on all 18 feature distributions. Alert if score exceeds 0.2. Auto-notification to ML team."),
         ("RETRAINING CADENCE","Models retrained every 90 days on latest labelled underwriter outcomes. Accuracy gate enforced before deployment."),
         ("AUDIT LOGS","Every decision logged: inputs, all agent outputs, final routing, confidence score, and timestamp — immutable."),
         ("COMPLIANCE","IRDAI-aligned explainability. Every rejection includes a written reason code. SHAP stored per high-value case.")]
for i,(title,d) in enumerate(pillars):
    r,c=divmod(i,2); x=0.4+c*6.5; y=1.05+r*2.5
    box(S18,x,y,6.1,2.3, fc=D3 if (r+c)%2==0 else D2, lc=MID, lw=Pt(0.75))
    tx(S18,title,x+0.15,y+0.15,5.8,0.4, sz=Pt(11), bold=True, color=WH)
    tx(S18,d,x+0.15,y+0.65,5.8,1.45, sz=Pt(12), color=LG)
bullets(S18,[
    "No model deployed without 90%+ accuracy on held-out validation — enforced automatically",
    "Feedback loop: underwriter escalation outcomes flow back into next training cycle",
],0.5,6.1,12.0)

S19 = sl(); hdr(S19,19,"Business Impact")
for i,(v,lbl) in enumerate([("<2 sec","Decision Turnaround"),("+34%","Conversion Lift"),("91%","Fraud Catch Rate"),("−62%","Manual Review Volume"),("4.8/5","Agent NPS Score")]):
    x=0.4+i*2.6
    box(S19,x,1.0,2.4,1.8, fc=D3, lc=MID, lw=Pt(0.5))
    tx(S19,v,x,1.1,2.4,0.85, sz=Pt(28), bold=True, color=WH, align=PP_ALIGN.CENTER)
    tx(S19,lbl,x,1.95,2.4,0.7, sz=Pt(10), color=DIM, align=PP_ALIGN.CENTER)
for i,(lbl,d) in enumerate([
    ("Intelligent Routing","HIGH / MED / LOW decisions with full audit — underwriter effort focused where it  matters most"),
    ("Premium Accuracy","Right price for every risk segment — eliminates both under and over pricing at source"),
    ("Customer Experience","Personalised plan with matched add-ons — reduces quote abandonment and improves NPS"),
]):
    y=3.1+i*1.2
    box(S19,0.4,y,12.5,1.05, fc=D3 if i%2==0 else D2, lc=MID, lw=Pt(0.5))
    tx(S19,lbl,0.6,y+0.25,3.8,0.55, sz=Pt(13), bold=True, color=WH)
    tx(S19,d,4.6,y+0.28,8.2,0.5, sz=Pt(12), color=LG)

S20 = sl(); hdr(S20,20,"Roadmap & Call to Action")
phases=[("Q2 2026","PILOT","Deploy on 500 live quotes per week. Measure decision accuracy vs manual baseline."),
        ("Q3 2026","SCALE","Expand to all new business lines. Integrate real-time claims feed for fraud retraining."),
        ("Q4 2026","OPTIMISE","AutoML model selection. A/B test personalisation strategies across customer segments."),
        ("2027+","EVOLVE","LLM-assisted underwriting notes. Predictive lifetime value scoring at quote time.")]
for i,(q,phase,d) in enumerate(phases):
    x=0.4+i*3.25
    box(S20,x,1.05,3.05,2.8, fc=D3 if i%2==0 else D2, lc=MID, lw=Pt(0.75))
    tx(S20,q,x+0.15,1.15,2.75,0.35, sz=Pt(8), color=DIM)
    tx(S20,phase,x+0.15,1.5,2.75,0.6, sz=Pt(18), bold=True, color=WH)
    tx(S20,d,x+0.15,2.15,2.75,1.5, sz=Pt(11), color=LG)
box(S20,0.4,4.1,12.5,0.04, fc=MID)
tx(S20,"YOUR NEXT STEPS",0.6,4.3,12.0,0.5, sz=Pt(16), bold=True, color=WH, align=PP_ALIGN.CENTER)
for i,a in enumerate([
    "✓  Approve pilot scope and data access for Q2 2026 deployment",
    "✓  Schedule 2-week sprint review with AI and underwriting teams",
    "✓  Set baseline KPIs to measure against in the 90-day pilot report",
]):
    tx(S20,a,1.5,4.95+i*0.6,10.5,0.55, sz=Pt(14), color=LG, align=PP_ALIGN.CENTER)

S21 = sl(); hdr(S21,21,"Appendix — Optional Deep-Dive Slides")
for i,(ref,title,d) in enumerate([
    ("A1","Data Dictionary","Field-by-field description of all 23 input features — type, valid range, canonical Enum string, and source"),
    ("A2","Model Metrics Card","Train/test accuracy, precision, recall, F1, and AUC for all five models — pre-deployment gate results"),
    ("A3","Governance Checklist","IRDAI compliance checklist, audit trail requirements, model risk management, disaster recovery"),
    ("A4","API Documentation","OpenAPI spec for /evaluate_quote endpoint — sample request, full JSON response schema"),
    ("A5","Taxonomy Reference","Canonical Enum strings for all categorical fields — aligned across UI, API, and model encoders"),
    ("A6","Agent Error Handling","Fallback logic if an agent fails mid-pipeline — graceful degradation design and ops alerting"),
]):
    y=1.05+i*0.98
    box(S21,0.4,y,12.5,0.88, fc=D3 if i%2==0 else D2, lc=MID, lw=Pt(0.5))
    tx(S21,ref,0.6,y+0.22,0.7,0.5, sz=Pt(11), bold=True, color=WH)
    tx(S21,title,1.5,y+0.22,3.5,0.45, sz=Pt(13), bold=True, color=WH)
    tx(S21,d,5.3,y+0.25,7.5,0.42, sz=Pt(11), color=LG)

OUT = r"C:\Users\mamid\Desktop\QuoteAI_Cinematic.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
